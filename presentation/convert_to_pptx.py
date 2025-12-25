"""
HTML Presentation to PPTX Converter
===================================
이 스크립트는 presentation_en.html의 내용을 PPTX로 변환합니다.

사용법:
    pip install python-pptx Pillow
    python convert_to_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

def rgb(r, g, b):
    """RGB 색상을 16진수 문자열로 변환"""
    return f"{r:02X}{g:02X}{b:02X}"

# Colors (hex strings)
AIRBNB_RED = rgb(255, 56, 92)
BLACK = rgb(0, 0, 0)
WHITE = rgb(255, 255, 255)
GRAY_100 = rgb(247, 247, 247)
GRAY_500 = rgb(113, 113, 113)
GRAY_700 = rgb(72, 72, 72)
TEAL = rgb(0, 166, 153)

def set_shape_color(shape, color_hex):
    """도형 채우기 색상 설정"""
    from pptx.dml.color import RGBColor
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16))

def set_font_color(font, color_hex):
    """폰트 색상 설정"""
    from pptx.dml.color import RGBColor
    font.color.rgb = RGBColor(int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16))

def add_title_slide(prs, title, subtitle="", is_dark=False):
    """타이틀 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(background, BLACK if is_dark else WHITE)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    set_font_color(p.font, WHITE if is_dark else BLACK)
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        set_font_color(p2.font, GRAY_500)
        p2.space_before = Pt(20)
    
    return slide

def add_content_slide(prs, title, content_items, is_dark=False):
    """컨텐츠 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(bg, BLACK if is_dark else WHITE)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    set_font_color(p.font, WHITE if is_dark else BLACK)
    
    # Content
    y_pos = 1.8
    for item in content_items:
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11), Inches(0.6))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {item}" if item else ""
        p.font.size = Pt(20)
        set_font_color(p.font, WHITE if is_dark else GRAY_700)
        y_pos += 0.6
    
    return slide

def add_metric_card(slide, x, y, label, value, accent=False):
    """메트릭 카드 추가"""
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(1.3))
    set_shape_color(card, rgb(255, 245, 247) if accent else GRAY_100)
    card.line.fill.background()
    
    # Label
    label_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(2.2), Inches(0.3))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    set_font_color(p.font, GRAY_500)
    
    # Value
    value_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.5), Inches(2.2), Inches(0.6))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(32)
    p.font.bold = True
    set_font_color(p.font, AIRBNB_RED if accent else BLACK)

def create_presentation():
    """프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = add_title_slide(prs, "Airbnb\nSuperhost\nPrediction", 
                           "Building ML Classification Models to Predict Airbnb Superhosts", 
                           is_dark=True)
    
    # Slide 2: Problem Definition
    slide = add_content_slide(prs, "Why Predict Superhost Status?", [
        "22% Revenue Increase — Superhosts earn more than regular hosts",
        "60% Higher Booking Rate — Badge increases guest trust",
        "$100 Travel Coupon — Annual benefit for Superhosts",
        "",
        "Core Question: Can we predict which hosts will become Superhosts?"
    ])
    
    # Slide 3: Data Overview
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(bg, GRAY_100)
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Dataset Structure"
    p.font.size = Pt(44)
    p.font.bold = True
    
    add_metric_card(slide, 0.8, 2, "TRAINING SAMPLES", "25,386")
    add_metric_card(slide, 3.5, 2, "TEST SAMPLES", "130")
    add_metric_card(slide, 0.8, 3.5, "ORIGINAL FEATURES", "54")
    add_metric_card(slide, 3.5, 3.5, "FINAL FEATURES", "444", accent=True)
    
    # Class distribution text
    dist_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5), Inches(3))
    tf = dist_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Class Distribution"
    p.font.size = Pt(24)
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = "\n• Not Superhost: 64%\n• Superhost: 36%\n\n* Upsampling applied"
    p2.font.size = Pt(18)
    set_font_color(p2.font, GRAY_700)
    
    # Slide 4: Preprocessing
    slide = add_content_slide(prs, "Data Preprocessing Pipeline", [
        "Step 1: Missing Values — Numeric: Median, Categorical: Mode",
        "Step 2: Feature Engineering — Date → days_since, Price/% → Numeric",
        "Step 3: Text Embedding — Amenities → 384-dim (SentenceTransformer)",
        "Step 4: Encoding & Scaling — One-Hot + StandardScaler → 444 Features",
        "",
        "Property Type: 76 types → 6 groups via embedding clustering",
        "Class Balancing: 13,731:6,577 → 13,731:13,731 (Upsampling)"
    ])
    
    # Slide 5: Amenity Embedding
    slide = add_content_slide(prs, "Amenity Feature Embedding", [
        "Problem: 6,372 unique amenity terms → sparse one-hot encoding",
        "Solution: SentenceTransformer (all-MiniLM-L6-v2) → 384-dim dense vector",
        "",
        "Benefits:",
        "  • Captures semantic similarity between amenities",
        "  • Reduces dimensionality: 6,372 → 384",
        "  • Handles unseen amenity combinations",
        "",
        "Categories: Entertainment/Tech, Kitchen/Dining, Bathroom/Bedroom,",
        "Climate Control, Safety, Facilities/Outdoors, Access/Parking, Other"
    ])
    
    # Slide 6: Model Comparison
    slide = add_content_slide(prs, "Model Performance Comparison", [
        "Random Forest      — Acc: 93.0% | Recall: 67.5% | F1: 73.6%  ✓ Selected",
        "Decision Tree      — Acc: 89.4% | Recall: 66.2% | F1: 67.8%",
        "PyTorch MLP        — Acc: 83.5% | Recall: 82.7% | F1: 76.4%  (Best F1)",
        "KNN                — Acc: 82.9% | Recall: 67.9% | F1: 60.8%",
        "Logistic Regression — Acc: 80.2% | Recall: 82.5% | F1: 72.5%",
        "Naive Bayes        — Acc: 72.1% | Recall: 87.1% | F1: 58.2%  (Best Recall)"
    ], is_dark=True)
    
    # Slide 7: Model Selection
    slide = add_content_slide(prs, "Why Random Forest?", [
        "Selected Model: Random Forest (Acc: 93.0%, Recall: 67.5%, F1: 73.6%)",
        "",
        "Selection Rationale:",
        "  • High Accuracy — Reliable overall predictions",
        "  • Feature Importance — Interpretable model",
        "  • Stability — Low overfitting risk",
        "  • Business View — Avoid false Superhost awards",
        "",
        "Top Features: host_listings_count, review_scores_rating,",
        "number_of_reviews_ltm, host_response_rate, review_scores_value"
    ])
    
    # Slide 8: Airbnb Criteria
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(bg, AIRBNB_RED)
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Airbnb Official\nSuperhost Criteria"
    p.font.size = Pt(44)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    
    criteria = [
        ("10+", "Completed stays/year"),
        ("90%", "Response Rate"),
        ("<1%", "Cancellation Rate"),
        ("4.8", "Overall Rating")
    ]
    
    x_start = 0.8
    for i, (val, label) in enumerate(criteria):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                     Inches(x_start + i * 3.1), Inches(3), 
                                     Inches(2.8), Inches(2))
        set_shape_color(box, rgb(200, 50, 80))
        box.line.fill.background()
        
        val_box = slide.shapes.add_textbox(Inches(x_start + i * 3.1 + 0.2), Inches(3.2), Inches(2.4), Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(40)
        p.font.bold = True
        set_font_color(p.font, WHITE)
        
        label_box = slide.shapes.add_textbox(Inches(x_start + i * 3.1 + 0.2), Inches(4.2), Inches(2.4), Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        set_font_color(p.font, WHITE)
    
    # Slide 9: Comparison
    slide = add_content_slide(prs, "ML Model vs Official Criteria", [
        "Airbnb Official:                    ML Model Top Features:",
        "• Response Rate ≥ 90%              • host_listings_count",
        "• Rating ≥ 4.8                     • review_scores_rating",
        "• Annual Reviews ≥ 10              • number_of_reviews_ltm",
        "• Cancellation < 1%                • host_response_rate",
        "",
        "✓ Aligned: Response Rate, Rating, Activity metrics",
        "✕ ML-only: host_listings_count, amenities embedding, experience duration"
    ])
    
    # Slide 10: Disagreement Analysis
    slide = add_content_slide(prs, "Disagreement Analysis", [
        "Agreement Rate: 80.8% | Disagreement: 25 cases (19.2%)",
        "",
        "Case #3: Response 100% ✓, Rating 4.90 ✓, Reviews 18 ✓ → Rule: ✓, ML: ✕",
        "Case #16: Response 100% ✓, Rating 5.00 ✓, Reviews 10 ✓ → Mystery case",
        "Case #5: Response 100% ✓, Rating 5.00 ✓, Reviews 8 ✕ → Under 10 reviews",
        "",
        "Key Insight: ML model learns latent patterns like 'hosting scale'",
        "and 'experience duration' not in official criteria"
    ], is_dark=False)
    
    # Slide 11: Conclusion
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(bg, BLACK)
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Conclusions & Implications"
    p.font.size = Pt(44)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    
    conclusions = [
        ("High Accuracy", "Random Forest: 93%"),
        ("Criteria Alignment", "Agreement: 80.8%"),
        ("Hidden Patterns", "Hosting scale & experience")
    ]
    
    for i, (title, desc) in enumerate(conclusions):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.8 + i * 4), Inches(2),
                                     Inches(3.5), Inches(2))
        set_shape_color(box, GRAY_700)
        box.line.fill.background()
        
        t_box = slide.shapes.add_textbox(Inches(1 + i * 4), Inches(2.3), Inches(3.2), Inches(1.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        set_font_color(p.font, WHITE)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(24)
        set_font_color(p2.font, WHITE)
    
    # Applications box
    app_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2))
    set_shape_color(app_box, AIRBNB_RED)
    app_box.line.fill.background()
    
    app_text = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.3), Inches(1.6))
    tf = app_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Applications"
    p.font.size = Pt(20)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    p2 = tf.add_paragraph()
    p2.text = "Host Coaching: Identify high-potential hosts early | Platform Operations: Quality management & marketing optimization"
    p2.font.size = Pt(16)
    set_font_color(p2.font, WHITE)
    
    return prs

if __name__ == "__main__":
    print("Creating PPTX presentation...")
    prs = create_presentation()
    
    output_path = "presentation_en.pptx"
    prs.save(output_path)
    print(f"✅ Saved to: {output_path}")
    print(f"   File size: {os.path.getsize(output_path) / 1024:.1f} KB")

