"""
HTML Presentation to PPTX Converter (Korean)
==========================================
이 스크립트는 presentation_ko.html의 내용을 PPTX로 변환합니다.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def rgb(r, g, b):
    """RGB 색상을 16진수 문자열로 변환"""
    return f"{r:02X}{g:02X}{b:02X}"

# Colors (hex strings)
AIRBNB_RED = rgb(255, 56, 92)
BLACK = rgb(0, 0, 0)
WHITE = rgb(255, 255, 255)
GRAY_100 = rgb(247, 247, 247)
GRAY_500 = rgb(113, 113, 113)
GRAY_700 = rgb(72, 72, 72)

def get_rgb_color(color_hex):
    from pptx.dml.color import RGBColor
    return RGBColor(int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16))

def set_shape_color(shape, color_hex):
    """도형 채우기 색상 설정"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = get_rgb_color(color_hex)

def set_font_color(font, color_hex):
    """폰트 색상 설정"""
    font.color.rgb = get_rgb_color(color_hex)

def add_title_slide(prs, title, subtitle="", is_dark=False):
    """타이틀 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(bg, BLACK if is_dark else WHITE)
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11), Inches(3))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    set_font_color(p.font, WHITE if is_dark else BLACK)
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        set_font_color(p2.font, GRAY_500)
        p2.space_before = Pt(20)
    
    return slide

def add_content_slide(prs, title, content_items, is_dark=False):
    """컨텐츠 슬라이드 추가"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(bg, BLACK if is_dark else WHITE)
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    set_font_color(p.font, WHITE if is_dark else BLACK)
    
    y_pos = 1.8
    for item in content_items:
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11), Inches(0.6))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• {item}" if item else ""
        p.font.size = Pt(20)
        set_font_color(p.font, WHITE if is_dark else GRAY_700)
        y_pos += 0.6
    
    return slide

def add_metric_card(slide, x, y, label, value, accent=False):
    """메트릭 카드 추가"""
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(1.3))
    set_shape_color(card, rgb(255, 245, 247) if accent else GRAY_100)
    card.line.fill.background()
    
    label_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(2.2), Inches(0.3))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    set_font_color(p.font, GRAY_500)
    
    value_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.5), Inches(2.2), Inches(0.6))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(32)
    p.font.bold = True
    set_font_color(p.font, AIRBNB_RED if accent else BLACK)

def create_presentation():
    """프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, "Airbnb\nSuperhost\nPrediction", 
                   "에어비앤비 슈퍼호스트 예측을 위한 머신러닝 모델 구축", 
                   is_dark=True)
    
    # Slide 2: Problem Definition
    add_content_slide(prs, "슈퍼호스트 예측이 왜 중요한가?", [
        "22% 수익 증가 — 슈퍼호스트는 일반 호스트보다 높은 수익 창출",
        "60% 높은 예약률 — 뱃지가 게스트의 신뢰도를 높임",
        "$100 여행 쿠폰 — 슈퍼호스트를 위한 연간 혜택",
        "",
        "핵심 질문: 어떤 호스트가 슈퍼호스트가 될지 미리 알 수 있을까?"
    ])
    
    # Slide 3: Data Overview
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(bg, GRAY_100)
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "데이터셋 구조"
    p.font.size = Pt(44)
    p.font.bold = True
    
    add_metric_card(slide, 0.8, 2, "학습 데이터", "25,386")
    add_metric_card(slide, 3.5, 2, "테스트 데이터", "130")
    add_metric_card(slide, 0.8, 3.5, "원본 피처", "54")
    add_metric_card(slide, 3.5, 3.5, "최종 피처", "444", accent=True)
    
    dist_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5), Inches(3))
    tf = dist_box.text_frame
    p = tf.paragraphs[0]
    p.text = "클래스 분포"
    p.font.size = Pt(24)
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = "\n• 일반 호스트: 64%\n• 슈퍼호스트: 36%\n\n* Upsampling 적용됨"
    p2.font.size = Pt(18)
    set_font_color(p2.font, GRAY_700)
    
    # Slide 4: Preprocessing
    add_content_slide(prs, "데이터 전처리 파이프라인", [
        "1단계: 결측치 처리 — 수치형: 중앙값, 범주형: 최빈값",
        "2단계: 피처 엔지니어링 — 날짜 → 경과일수, 가격/% → 수치형 변환",
        "3단계: 텍스트 임베딩 — 편의시설(Amenities) → 384차원 벡터",
        "4단계: 인코딩 & 스케일링 — One-Hot + StandardScaler",
        "",
        "숙소 유형: 76개 유형 → 임베딩 클러스터링으로 6개 그룹화",
        "클래스 균형: 13,731:6,577 → 13,731:13,731 (Upsampling)"
    ])
    
    # Slide 5: Amenity Embedding (Updated Layout)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(bg, WHITE)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "어메니티 임베딩"
    p.font.size = Pt(44)
    p.font.bold = True
    
    # Model Info
    info_box = slide.shapes.add_textbox(Inches(9), Inches(0.8), Inches(3.5), Inches(0.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Model: SentenceTransformer (all-MiniLM-L6-v2)"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.RIGHT
    set_font_color(p.font, GRAY_500)
    
    # Left Column: Description
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(2))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "원본 데이터에는 6,372개의 고유한 편의시설 용어가 존재합니다. 이를 전통적인 One-Hot 인코딩으로 처리하면 데이터가 매우 희소해져 의미적 연결성을 잃게 됩니다."
    p.font.size = Pt(16)
    set_font_color(p.font, GRAY_700)
    p.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = "우리는 이 용어들을 384차원 밀집 벡터로 임베딩하여, 'Wifi'와 'Internet' 같은 단어 간의 의미적 유사성을 포착했습니다."
    p2.font.size = Pt(16)
    set_font_color(p2.font, GRAY_700)
    
    # Process Flow (Simplified Text)
    flow_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4), Inches(5.5), Inches(0.8))
    set_shape_color(flow_box, GRAY_100)
    flow_box.line.fill.background()
    
    flow_text = slide.shapes.add_textbox(Inches(0.9), Inches(4.1), Inches(5.3), Inches(0.6))
    tf = flow_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📝 6,372 용어  →  🤖 Transformer  →  📊 384차원"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    
    # Benefits Box
    ben_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.2), Inches(5.5), Inches(1.2))
    set_shape_color(ben_box, rgb(255, 245, 247))
    ben_box.line.color.rgb = get_rgb_color(AIRBNB_RED)
    
    ben_text = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(5.1), Inches(1))
    tf = ben_text.text_frame
    p = tf.paragraphs[0]
    p.text = "단순한 유무(Presence)를 넘어, 모델이 편의시설의 품질과 완결성을 이해하도록 돕습니다."
    p.font.size = Pt(14)
    p.font.bold = True
    set_font_color(p.font, AIRBNB_RED)
    
    # Right Column: Placeholder for Map
    map_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.6))
    set_shape_color(map_box, GRAY_100)
    map_box.line.color.rgb = get_rgb_color(GRAY_500)
    map_box.line.dash_style = 1 # Solid
    
    map_label = slide.shapes.add_textbox(Inches(7), Inches(3.5), Inches(5.3), Inches(1))
    tf = map_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Interactive t-SNE Visualization\n(HTML 파일 참조)"
    p.alignment = PP_ALIGN.CENTER
    set_font_color(p.font, GRAY_500)

    # Slide 6: Model Comparison
    add_content_slide(prs, "모델 성능 비교", [
        "Random Forest      — 정확도: 93.0% | 재현율: 67.5% | F1: 73.6%  ✓ 선택됨",
        "Decision Tree      — 정확도: 89.4% | 재현율: 66.2% | F1: 67.8%",
        "PyTorch MLP        — 정확도: 83.5% | 재현율: 82.7% | F1: 76.4%  (최고 F1)",
        "KNN                — 정확도: 82.9% | 재현율: 67.9% | F1: 60.8%",
        "Logistic Regression — 정확도: 80.2% | 재현율: 82.5% | F1: 72.5%",
        "Naive Bayes        — 정확도: 72.1% | 재현율: 87.1% | F1: 58.2%  (최고 재현율)"
    ], is_dark=True)
    
    # Slide 7: Model Selection
    add_content_slide(prs, "왜 Random Forest인가?", [
        "선택 모델: Random Forest (정확도: 93.0%, F1: 73.6%)",
        "",
        "선택 이유:",
        "  • 높은 정확도 — 전반적으로 가장 신뢰할 수 있는 예측",
        "  • 변수 중요도 — 모델의 판단 근거 해석 용이",
        "  • 안정성 — 과적합 위험이 낮음",
        "  • 비즈니스 관점 — 잘못된 슈퍼호스트 선정 방지 (높은 정밀도)",
        "",
        "상위 중요 변수: 호스트 리스팅 수, 평점, 최근 리뷰 수, 응답률"
    ])
    
    # Slide 8: Airbnb Criteria
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(bg, AIRBNB_RED)
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "에어비앤비 공식\n슈퍼호스트 기준"
    p.font.size = Pt(44)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    
    criteria = [
        ("10+", "연간 완료된 숙박"),
        ("90%", "응답률"),
        ("<1%", "취소율"),
        ("4.8", "전체 평점")
    ]
    
    x_start = 0.8
    for i, (val, label) in enumerate(criteria):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                     Inches(x_start + i * 3.1), Inches(3), 
                                     Inches(2.8), Inches(2))
        set_shape_color(box, rgb(200, 50, 80))
        box.line.fill.background()
        
        val_box = slide.shapes.add_textbox(Inches(x_start + i * 3.1 + 0.2), Inches(3.2), Inches(2.4), Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(40)
        p.font.bold = True
        set_font_color(p.font, WHITE)
        
        label_box = slide.shapes.add_textbox(Inches(x_start + i * 3.1 + 0.2), Inches(4.2), Inches(2.4), Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        set_font_color(p.font, WHITE)
    
    # Slide 9: Comparison
    add_content_slide(prs, "ML 모델 vs 공식 기준", [
        "에어비앤비 공식 기준:                ML 모델 상위 중요 변수:",
        "• 응답률 ≥ 90%                     • host_listings_count (리스팅 수)",
        "• 평점 ≥ 4.8                       • review_scores_rating (평점)",
        "• 연간 숙박 수 ≥ 10                • number_of_reviews_ltm (최근 리뷰)",
        "• 취소율 < 1%                      • host_response_rate (응답률)",
        "",
        "✓ 일치하는 점: 응답률, 평점, 활동성 지표",
        "✕ ML만의 발견: 호스트 리스팅 수(규모), 어메니티 품질, 경력 기간"
    ])
    
    # Slide 10: Disagreement Analysis
    add_content_slide(prs, "불일치 케이스 분석", [
        "일치율: 80.8% | 불일치: 25건 (19.2%)",
        "",
        "케이스 #3: 응답 100% ✓, 평점 4.90 ✓, 리뷰 18 ✓ → 규칙: ✓, ML: ✕",
        "케이스 #16: 응답 100% ✓, 평점 5.00 ✓, 리뷰 10 ✓ → 미스터리 케이스",
        "케이스 #5: 응답 100% ✓, 평점 5.00 ✓, 리뷰 8 ✕ → 리뷰 수 부족",
        "",
        "핵심 인사이트: ML 모델은 공식 기준에 없는 '운영 규모'나",
        "'경력 기간' 같은 잠재적 패턴을 학습하여 더 유연하게 판단함"
    ], is_dark=False)
    
    # Slide 11: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    set_shape_color(bg, BLACK)
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "결론 및 시사점"
    p.font.size = Pt(44)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    
    conclusions = [
        ("높은 정확도", "Random Forest: 93%"),
        ("기준 일치도", "일치율: 80.8%"),
        ("숨겨진 패턴", "운영 규모 & 경력")
    ]
    
    for i, (title, desc) in enumerate(conclusions):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(0.8 + i * 4), Inches(2),
                                     Inches(3.5), Inches(2))
        set_shape_color(box, GRAY_700)
        box.line.fill.background()
        
        t_box = slide.shapes.add_textbox(Inches(1 + i * 4), Inches(2.3), Inches(3.2), Inches(1.5))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        set_font_color(p.font, WHITE)
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(24)
        set_font_color(p2.font, WHITE)
    
    # Applications box
    app_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2))
    set_shape_color(app_box, AIRBNB_RED)
    app_box.line.fill.background()
    
    app_text = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.3), Inches(1.6))
    tf = app_text.text_frame
    p = tf.paragraphs[0]
    p.text = "활용 방안"
    p.font.size = Pt(20)
    p.font.bold = True
    set_font_color(p.font, WHITE)
    p2 = tf.add_paragraph()
    p2.text = "호스트 코칭: 잠재력 높은 호스트 조기 발굴 | 플랫폼 운영: 품질 관리 및 마케팅 최적화"
    p2.font.size = Pt(16)
    set_font_color(p2.font, WHITE)
    
    return prs

if __name__ == "__main__":
    print("Creating PPTX presentation (Korean)...")
    prs = create_presentation()
    output_path = "presentation_ko.pptx"
    prs.save(output_path)
    print(f"✅ Saved to: {output_path}")
    print(f"   File size: {os.path.getsize(output_path) / 1024:.1f} KB")
